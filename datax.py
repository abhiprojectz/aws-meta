from PIL import Image, ImageDraw, ImageFont
import string
import random
import csv
import requests 
import shutil
from time import sleep
import subprocess
import textwrap 
from turtle import width
from pptx import Presentation
from pptx.util import Inches
import os
import multiprocessing
import subprocess
from time import sleep
from pyvirtualdisplay import Display


def id_generator(size=20, chars=string.ascii_uppercase + string.digits):
  return ''.join(random.choice(chars) for _ in range(size))

# __file__ = "/content/dddfff.py"

#circleCi
platform_path = "/home/circleci/project/"
# Colab
# platform_path = "/content/"

dir_path = os.path.dirname(os.path.realpath(__file__))
fps = 24


def downloadImg(imgs):
    # imgs = getImage()
    heads = {'User-Agent': 'Mozilla/5.0 (X11; Linux x86_64) AppleWebKit/537.11 (KHTML, like Gecko) Chrome/23.0.1271.64 Safari/537.11',
       'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8',
       'Accept-Charset': 'ISO-8859-1,utf-8;q=0.7,*;q=0.3',
       'Accept-Encoding': 'none',
       'Accept-Language': 'en-US,en;q=0.8',
       'Connection': 'keep-alive'}
    for i in imgs:
        dir_path = os.path.dirname(os.path.realpath(__file__))
        image_url = i
        filename = "card_" + str(imgs.index(i)) + ".jpg"
        fz = os.path.join(dir_path, "imgs")
        r = requests.get(image_url,headers=heads, stream = True)
        if r.status_code == 200:
            r.raw.decode_content = True
            with open(filename,'wb') as f:
                shutil.copyfileobj(r.raw, f)

            shutil.move(str(filename), fz)
            print('Image sucessfully Downloaded: ', filename)
            sleep(2)
        else:
            print('Image Couldn\'t be retreived')

def get_data(filename): 
  # filename = "/content/data.csv"
  item1 = []
  item2 = []
  # item3 = []

  with open(filename, 'r') as csvfile:
      csvreader = csv.reader(csvfile)
      fields = next(csvreader)
      for row in csvreader:
          item1.append(row[0])
          item2.append(row[2])
          # item3.append(row[3])
  return item1, item2

def generate_card(num , main_text, main_text2, main_img):
  out_png = os.path.join(dir_path, "generated/" + "card_" + str(num) + '.png')
  temp_bg = Image.open(os.path.join(dir_path, "temp_bg.PNG")) 
  img_main = Image.open(os.path.join(dir_path, "imgs/" + main_img)) 
  # x_size = (500, 720)
  # black_bg = Image.new('RGB', x_size)
  main_W = temp_bg.width
  main_H = temp_bg.height

  ff = os.path.join(dir_path, 'Raleway-Black.ttf')
  I1 = ImageDraw.Draw(temp_bg)

  font_num = ImageFont.truetype( ff ,60)
  w1, h1 = I1.textsize(num, font_num)
  I1.text(((main_W - w1)/2, 50), num , fill=(255, 255, 255), align="center", font=font_num)
  
  msg = textwrap.fill(main_text, 20) 
  font_title = ImageFont.truetype(ff, 30)
  w2, h2 = I1.textsize(msg, font_title)
  h3 = 0
  if main_text2:
    msg2 = main_text2
    font_title2 = ImageFont.truetype(ff, 18)
    w3, h3 = I1.textsize(msg2, font_title2)
  I1.text(((main_W - w2)/2, (100-h2)/2 + 220 - h3), msg, fill=(0, 0, 0),align="center", font=font_title)


  if main_text2:
    msg2 = main_text2
    font_title2 = ImageFont.truetype(ff, 18)
    w3, h3 = I1.textsize(msg2, font_title2)
    I1.text(((main_W - w3)/2, (100 + h2 -h3)/2 + 230), msg2, fill=(0, 0, 0), align="center", font=font_title2)


  img_main_size = (300,300)
  back_im = temp_bg.copy()
  back_im.paste(img_main.resize(img_main_size),  (50, 400))
  back_im.save(out_png)

# num = "7"
# main_text = "Teenage Mutant loria tolas "
# main_text2 = "Golam 7778$"
# main_img = "card_1.jpg"
# generate_card(num ,main_text , main_text2, main_img) 

def get_final_img_paths():
    files = []
    _downloads_path = os.path.join(dir_path, "imgs/")

    for file in os.listdir(_downloads_path):
        if file.endswith(".jpg"):
            # tp = os.path.join(fz, file)
            tp = file
            files.append(tp)
            # print(os.path.join(file))
    return files

def gen_black():
  out_png1 = os.path.join("card_black_1.png")
  out_png2 = os.path.join("card_black_2.png")
  out_png3 = os.path.join("card_black_3.png")
  temp_bg = Image.open(os.path.join(dir_path, "temp_bg.PNG")) 
  main_W = temp_bg.width
  main_H = temp_bg.height
  black_bg = Image.new('RGB', (main_W, main_H))
  black_bg.save(out_png1)
  black_bg.save(out_png2)
  black_bg.save(out_png3)

def get_bg_music():
    _bg_m_path = os.path.join(dir_path, "bg_music") 
    files = []
    for file in os.listdir(_bg_m_path):
        if file.endswith(".mp3"):
            files.append(os.path.join(dir_path, "bg_music")  + "/" + file)
    f = random.choice(files)
    return f

# OLD ====================
# def generate_vid():
#   os.environ['DURATION'] = '20'
#   _slides = len(get_final_img_paths())
#   _dur = str(int(4.3*_slides))
#   env = {'DURATION': _dur}
#   process = subprocess.call("bash gen_stack.sh",env=env, shell=True)
#   # sleep(2)
#   dest = os.path.join(dir_path, "outputs")
#   bg_music = get_bg_music()
#   print(bg_music)
#   _inp = "/content/outputs/stack_res.mp4"
#   _out = dest + "/stack_final.mp4"
#   process2 = subprocess.call(f"ffmpeg -y -i {_inp} -i {bg_music} -map 0:v -map 1:a -c:v copy -shortest {_out}", shell=True)

def generate_content(_names, _urls):
  # item1, item2, item3 =  get_data("/content/data.csv")
  item1 = _names
  item2 = _urls
  downloadImg(item2)
  imgs = get_final_img_paths()

  for i,j in zip(item1,item2):
    num = str(int(item1.index(i) + 1))
    # print(num)
    main_text = str(i)
    # main_text2 = "Market cap: " + str(j)
    main_text2 = None
    main_img = f'card_{item1.index(i)}.jpg'
    generate_card(num ,main_text , main_text2, main_img)
    print(f"Card {item1.index(i)} generated.")

def generate_vid():
  dest = os.path.join(dir_path, "outputs")
  bg_music = get_bg_music()
  _inp = f"{platform_path}outputs/stack_res.mp4"
  _out = dest + "/stack_final.mp4"
  process2 = subprocess.call(f"ffmpeg -y -i {_inp} -i {bg_music} -map 0:v -map 1:a -c:v copy -shortest {_out}", shell=True)

def remove_all():
    subprocess.call(f"rm {platform_path}outputs/*.mp4", shell=True)
    subprocess.call(f"rm {platform_path}imgs/*.jpg", shell=True)
    subprocess.call(f"rm {platform_path}generated/*.png", shell=True)


def generate_ppt():
    ppt = Presentation(dir_path + '/template_30.pptx')
    base_ = dir_path + "/generated/"
    slide = ppt.slides[0]
    shapes = slide.shapes

    for num in range(3, 33):
        _res = base_ + "card_" + str(num-2) + ".png" 
        slide.shapes[0].shapes.add_picture(_res, num*3557977, 0 , 3557977, 6858000)
    ppt.save(dir_path + '/template_30_generated.pptx')
    print("generated")

def start_xvfb():
    disp = Display(size=(1536, 864), color_depth=24)
    disp.start()
    print(os.environ['DISPLAY'])

def task_1():
  subprocess.run('sudo libreoffice --norestore --impress template_30_generated.pptx', shell=True)


def task_2():
  subprocess.run('sudo libreoffice --norestore --impress template_30_generated.pptx', shell=True)
  
def task_3(_path):
  subprocess.run(f'ffmpeg -f x11grab -y -framerate 30 -s 1536x864 -i :0.0 -c:v libx264 -preset superfast -crf 18 {_path}', shell=True)


def stop_(_time):
  sleep(_time)
  subprocess.run("killall ffmpeg", shell=True)


def yt_engine(_names, _urls, _time, _path):
    remove_all()
    sleep(2)
    generate_content(_names, _urls)
    generate_ppt()
    start_xvfb()
    sleep(3)
    p1 = multiprocessing.Process(target=task_1)
    p1.start()
    sleep(3)
    subprocess.run("killall soffice.bin", shell=True)
    sleep(3)
    p2 = multiprocessing.Process(target=task_2)
    p2.start()
    sleep(3)
    subprocess.run("xdotool key F5", shell=True)
    sleep(1)
    subprocess.run("xdotool key space", shell=True)
    p3 = multiprocessing.Process(target=task_3, args=_path)
    p3.start()
    stop_(_time)
    sleep(3)
    generate_vid()
    print("Vid generated!")
    subprocess.run("killall Xvfb", shell=True)


# if __name__ == "__main__":
#     yt_engine()